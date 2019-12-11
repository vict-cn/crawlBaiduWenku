import pptx
import os
from pptx.util import Inches

def pic_to_ppt(filename):  #前提是图片文件名全为数字，否则还需要修改
    if not os.path.exists(filename):
        os.mkdir(filename)
    
    ppt = pptx.Presentation()
    pic_path=[]  
    for i in os.walk(filename).__next__()[2]:
        if i.endswith('.png'):
            pic_path.append(i)
    #若是不全为数字，则可尝试运行下列代码
#    ls=[]
#    for png in pic_path:
#        s=''
#        for item in png:
#            if item<='9' and item>='0':
#                s+=item
#        ls.append(s+'.png')
#    pic_path=ls
    
    pic_path.sort(key=lambda item:int(item.split('.')[0]))
    for i in pic_path:
        i='{}/{}'.format(filename,i)
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.add_picture(i, Inches(0), Inches(0), Inches(10), Inches(7.5))
        
    fname='{}/{}.pptx'.format(filename,filename)
    ppt.save(fname)
    print('生成的文件在 {} 文件夹下的 {}.ppt 中'.format(filename,filename))

if __name__=='__main__':
    filename='德邦物流公司PPT'
    pic_to_ppt(filename)
    
    
    
    
    